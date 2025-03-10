import streamlit as st
import google.generativeai as genai
import pdfplumber
import pandas as pd
import plotly.express as px
import re
import concurrent.futures
import time
from PIL import Image
import pytesseract
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

# Set Tesseract path
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Configure Gemini AI
genai.configure(api_key="AIzaSyCCZ4KxaJ2JrDv1PL7n1_2rU_iq-4BB_7g")
model = genai.GenerativeModel("gemini-2.0-flash-lite")

# Streamlit UI Configuration
st.set_page_config(page_title="📊 Financial Report AI", layout="wide")

# Custom CSS for modern design
st.markdown(
    """
    <style>
    .stApp {
        background-color: #f5f5f5;
        color: #333333;
    }
    .stButton button {
        background-color: #4CAF50;
        color: white;
        border-radius: 5px;
        padding: 10px 20px;
        font-size: 16px;
    }
    .stButton button:hover {
        background-color: #45a049;
    }
    .stHeader {
        font-size: 24px;
        font-weight: bold;
        color: #4CAF50;
    }
    .stSubheader {
        font-size: 20px;
        font-weight: bold;
        color: #333333;
    }
    .stMarkdown {
        font-size: 16px;
        color: #555555;
    }
    .stProgress {
        background-color: #4CAF50;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# Landing Page
st.markdown(
    """
    <div style="text-align: center; padding: 40px; background: linear-gradient(135deg, #6a11cb, #2575fc); color: white; border-radius: 10px;">
        <h1 style="font-size: 48px; font-weight: bold;">📊 Financial Report AI</h1>
        <p style="font-size: 20px;">Your AI-powered tool for financial analysis and insights.</p>
    </div>
    """,
    unsafe_allow_html=True,
)

# Dark Mode Toggle
dark_mode = st.sidebar.toggle("Dark Mode", value=False)

# Apply custom CSS for dark mode
if dark_mode:
    st.markdown(
        """
        <style>
        .stApp {
            background-color: #1E1E1E;
            color: #FFFFFF;
        }
        .stHeader {
            color: #4CAF50;
        }
        .stSubheader {
            color: #FFFFFF;
        }
        .stMarkdown {
            color: #CCCCCC;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

# File Upload Section
st.sidebar.header("Upload Financial Report")
uploaded_file = st.sidebar.file_uploader("Choose a PDF or Image file", type=["pdf", "png", "jpg", "jpeg"])

# Summarization Length Control
st.sidebar.header("Summarization Length")
summarization_length = st.sidebar.radio("Choose summarization length:", ("Short", "Detailed"))

# Prediction Period Control
st.sidebar.header("Prediction Period")
prediction_period = st.sidebar.slider("Select Prediction Period (Years)", 1, 5, 1)

# --- FUNCTIONS ---
@st.cache_data
def extract_text_from_pdf(uploaded_file):
    """Extracts plain text from PDF pages."""
    with pdfplumber.open(uploaded_file) as pdf:
        text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
    return text

def extract_text_from_image(uploaded_file):
    """Extracts text from image using OCR."""
    image = Image.open(uploaded_file)
    text = pytesseract.image_to_string(image)
    return text

def extract_table_from_text(text):
    """Attempts to extract structured table data from raw extracted text."""
    lines = text.split("\n")
    structured_data = []

    for line in lines:
        values = re.findall(r"[\d,]+(?:\.\d+)?", line)  # Extract only numbers
        if values:
            structured_data.append(values)

    if structured_data:
        df = pd.DataFrame(structured_data)
        df = df.applymap(lambda x: float(str(x).replace(",", "")) if isinstance(x, (str, int, float)) else x)  # Clean numbers
        
        # Replace generic column names with actual dataset headings
        if df.shape[1] == 7:  # Ensure the number of columns matches the dataset
            df.columns = ["Time Period", "Revenue", "Expenses", "Profit", "Assets", "Liabilities", "Equity"]
        else:
            st.warning("The extracted data does not match the expected number of columns. Using generic names.")
            df.columns = [f"Metric {i+1}" for i in range(df.shape[1])]  # Fallback to generic names
        
        return df
    return None

def preprocess_data(df):
    """Cleans and preprocesses the data."""
    df = df.dropna()  # Drop rows with missing values
    df = df.drop_duplicates()  # Remove duplicate rows
    return df

def validate_data(df):
    """Validates the extracted data."""
    if df is None or df.empty:
        st.error("❌ No data found in the PDF. Please check the file and try again.")
        return False
    if len(df.select_dtypes(include=['number']).columns) < 2:
        st.error("❌ Insufficient numeric data for analysis. Please upload a valid financial report.")
        return False
    return True

def detect_time_series_column(df):
    """Finds the most likely column representing time periods."""
    df.columns = df.columns.astype(str)
    for col in df.columns:
        if any(keyword in col.lower() for keyword in ["date", "period", "time", "month", "year", "quarter"]):
            df[col] = pd.to_datetime(df[col], errors='coerce')
            return col
    return None

def generate_ai_insights(extracted_text, summarization_length):
    """Generates AI insights based on the extracted text and summarization length."""
    with st.spinner("Generating AI Insights..."):
        prompt = (
            f"Analyze this financial report and provide insights on Revenue, Expenses, Profit trends, and anomalies:\n"
            f"{extracted_text[:3000]}..."
        )
        if summarization_length == "Short":
            prompt += "\n\nProvide a concise summary."
        else:
            prompt += "\n\nProvide a detailed analysis."
        response = model.generate_content(prompt)
        return response.text

def generate_future_predictions(df, time_col, numeric_cols, prediction_period):
    """Generates predictions for the next year using Gemini AI."""
    with st.spinner("Generating Future Predictions..."):
        # If time_col is None, use the first numeric column as the x-axis
        if time_col is None:
            st.warning("No time-based column detected. Using the first numeric column for analysis.")
            columns_to_analyze = numeric_cols.tolist()
        else:
            columns_to_analyze = [time_col] + numeric_cols.tolist()
        
        # Create the prompt for Gemini AI
        prompt = (
            f"Based on the following financial data, predict the trends for the next {prediction_period} year(s):\n"
            f"{df[columns_to_analyze].to_string()}\n\n"
            f"Provide a detailed prediction for revenue, expenses, and profit."
        )
        
        # Generate the response using Gemini AI
        response = model.generate_content(prompt)
        return response.text

def analyze_sentiment(extracted_text):
    """Analyzes the sentiment of the financial report."""
    with st.spinner("Analyzing Sentiment..."):
        prompt = (
            f"Analyze the sentiment of the following financial report:\n"
            f"{extracted_text[:3000]}...\n\n"
            f"Is the sentiment positive, negative, or neutral? Provide a brief explanation."
        )
        response = model.generate_content(prompt)
        return response.text

def assess_risks(extracted_text):
    """Assesses potential risks in the financial report."""
    with st.spinner("Assessing Risks..."):
        prompt = (
            f"Identify potential risks in the following financial report:\n"
            f"{extracted_text[:3000]}...\n\n"
            f"Provide a detailed risk assessment."
        )
        response = model.generate_content(prompt)
        return response.text

def generate_recommendations(extracted_text):
    """Generates recommendations based on the financial report."""
    with st.spinner("Generating Recommendations..."):
        prompt = (
            f"Based on the following financial report, provide actionable recommendations:\n"
            f"{extracted_text[:3000]}...\n\n"
            f"Provide detailed recommendations for improving financial performance."
        )
        response = model.generate_content(prompt)
        return response.text

def plot_line_chart(df, x_axis, numeric_cols, dark_mode):
    """Plots a line chart with clear labels and annotations."""
    fig_line = px.line(df, x=x_axis, y=numeric_cols[1:], markers=True,
                       title="Revenue & Expenses Trend Over Time",
                       labels={x_axis: "Time Period", "value": "Amount (in USD)"},
                       hover_data={x_axis: "|%B %d, %Y"},
                       template="plotly_dark" if dark_mode else "plotly_white")
    
    # Add annotations for peaks and troughs
    for col in numeric_cols[1:]:
        max_value = df[col].max()
        min_value = df[col].min()
        max_index = df[col].idxmax()
        min_index = df[col].idxmin()
        
        fig_line.add_annotation(x=df[x_axis][max_index], y=max_value,
                                text=f"Max {col}: {max_value:,.2f}",
                                showarrow=True, arrowhead=1, ax=0, ay=-40)
        
        fig_line.add_annotation(x=df[x_axis][min_index], y=min_value,
                                text=f"Min {col}: {min_value:,.2f}",
                                showarrow=True, arrowhead=1, ax=0, ay=40)
    
    fig_line.update_traces(line=dict(width=2.5), selector=dict(mode='lines'))
    fig_line.update_layout(xaxis_title="Time Period", yaxis_title="Amount (in USD)",
                           legend_title="Metrics", hovermode="x unified")
    
    return fig_line

def plot_bar_chart(df, x_axis, numeric_cols, dark_mode):
    """Plots a bar chart with clear labels and annotations."""
    fig_bar = px.bar(df, x=x_axis, y=numeric_cols[1:], barmode="group",
                     title="Comparison of Financial Metrics",
                     labels={x_axis: "Time Period", "value": "Amount (in USD)"},
                     hover_data={x_axis: "|%B %d, %Y"},
                     template="plotly_dark" if dark_mode else "plotly_white")
    
    fig_bar.update_traces(marker=dict(line=dict(width=0.5, color='DarkSlateGrey')))
    fig_bar.update_layout(xaxis_title="Time Period", yaxis_title="Amount (in USD)",
                          legend_title="Metrics", hovermode="x unified")
    
    return fig_bar

def plot_scatter_chart(df, x_axis, numeric_cols, dark_mode):
    """Plots a scatter plot with clear labels and annotations."""
    fig_scatter = px.scatter(df, x=x_axis, y=numeric_cols[1:], 
                             title="Financial Metric Correlation",
                             labels={x_axis: "Time Period", "value": "Amount (in USD)"},
                             hover_data={x_axis: "|%B %d, %Y"},
                             template="plotly_dark" if dark_mode else "plotly_white")
    
    fig_scatter.update_traces(marker=dict(size=10, opacity=0.7))
    fig_scatter.update_layout(xaxis_title="Time Period", yaxis_title="Amount (in USD)",
                              legend_title="Metrics", hovermode="x unified")
    
    return fig_scatter

def create_presentation(df, ai_insights, sentiment, risks, recommendations):
    """Creates a PowerPoint presentation with the financial data and insights."""
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Financial Report Analysis"
    subtitle.text = "Generated by Financial Report AI"
    
    # Data Overview Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Overview"
    content = slide.shapes.placeholders[1].text_frame
    content.text = "Extracted Financial Data:"
    for index, row in df.iterrows():
        content.add_paragraph().text = str(row.to_dict())
    
    # AI Insights Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI Insights"
    content = slide.shapes.placeholders[1].text_frame
    content.text = ai_insights
    
    # Sentiment Analysis Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Sentiment Analysis"
    content = slide.shapes.placeholders[1].text_frame
    content.text = sentiment
    
    # Risk Assessment Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Risk Assessment"
    content = slide.shapes.placeholders[1].text_frame
    content.text = risks
    
    # Recommendations Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Recommendations"
    content = slide.shapes.placeholders[1].text_frame
    content.text = recommendations
    
    # Save the presentation to a BytesIO object
    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes

# --- MAIN LOGIC ---
if uploaded_file:
    with st.spinner("Extracting data from file..."):
        progress_bar = st.progress(0)
        for i in range(100):
            time.sleep(0.01)  # Simulate processing time
            progress_bar.progress(i + 1)
        
        if uploaded_file.type == "application/pdf":
            extracted_text = extract_text_from_pdf(uploaded_file)
        else:
            extracted_text = extract_text_from_image(uploaded_file)
        
        df = extract_table_from_text(extracted_text)

    if df is not None:
        df = preprocess_data(df)
        if not validate_data(df):
            st.stop()  # Stop execution if data is invalid

        st.subheader("📊 Extracted Financial Data")
        st.dataframe(df.style.set_properties({'text-align': 'center'}))

        # Display metric explanations
        st.markdown("---")
        st.subheader("📝 Metric Explanations")
        st.write("""
        - *Time Period*: The specific quarter and year for the financial data.
        - *Revenue*: Total income generated from business operations.
        - *Expenses*: Total costs incurred during business operations.
        - *Profit*: Net income calculated as Revenue minus Expenses.
        - *Assets*: Total resources owned by the company.
        - *Liabilities*: Total obligations or debts of the company.
        - *Equity*: The residual interest in the assets of the company after deducting liabilities.
        """)

        # Detect time series column
        time_col = detect_time_series_column(df)
        if time_col:
            df = df.sort_values(by=time_col)

        # Detect numeric columns
        numeric_cols = df.select_dtypes(include=['number']).columns

        if len(numeric_cols) > 1:
            x_axis = time_col if time_col else numeric_cols[0]

            # --- PLOTS ---
            st.markdown("---")
            st.subheader("📈 Financial Analysis Charts")

            col1, col2 = st.columns(2)
            with col1:
                # Line Chart with Tooltips and Smooth Animations
                fig_line = plot_line_chart(df, x_axis, numeric_cols, dark_mode)
                st.plotly_chart(fig_line, use_container_width=True)

            with col2:
                # Bar Chart with Tooltips and Animations
                fig_bar = plot_bar_chart(df, x_axis, numeric_cols, dark_mode)
                st.plotly_chart(fig_bar, use_container_width=True)

            st.subheader("📌 Data Distribution")
            # Scatter Plot with Tooltips and Animations
            fig_scatter = plot_scatter_chart(df, x_axis, numeric_cols, dark_mode)
            st.plotly_chart(fig_scatter, use_container_width=True)

            # --- AI INSIGHTS ---
            st.markdown("---")
            st.subheader("📄 AI-Generated Financial Summary")

            if extracted_text:
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(generate_ai_insights, extracted_text, summarization_length)
                    ai_insights = future.result()

                st.write(ai_insights)
                st.download_button("📥 Download AI Summary", ai_insights, file_name="Financial_Summary.txt")

            # --- SENTIMENT ANALYSIS ---
            st.markdown("---")
            st.subheader("📊 Sentiment Analysis")

            if extracted_text:
                sentiment = analyze_sentiment(extracted_text)
                st.write(sentiment)

            # --- RISK ASSESSMENT ---
            st.markdown("---")
            st.subheader("⚠ Risk Assessment")

            if extracted_text:
                risks = assess_risks(extracted_text)
                st.write(risks)

            # --- RECOMMENDATIONS ---
            st.markdown("---")
            st.subheader("📝 Recommendations")

            if extracted_text:
                recommendations = generate_recommendations(extracted_text)
                st.write(recommendations)

            # --- DOWNLOAD PROCESSED DATA ---
            st.markdown("---")
            st.subheader("📥 Download Processed Data")
            col1, col2 = st.columns(2)
            with col1:
                st.download_button("Download Processed Data (JSON)", df.to_json(orient="records"), file_name="financial_data.json")
            with col2:
                st.download_button("Download Processed Data (CSV)", df.to_csv(index=False), file_name="financial_data.csv")

            # --- CHATBOT FOR FUTURE PREDICTIONS ---
            st.markdown("---")
            st.subheader("🤖 Financial Predictions Chatbot")

            if st.button("Predict Future Trends"):
                predictions = generate_future_predictions(df, time_col, numeric_cols, prediction_period)
                st.write("*Chatbot:* Here are the predicted trends for the next year:")
                st.write(predictions)

            # --- DOWNLOAD PPTX ---
            st.markdown("---")
            st.subheader("📊 Download PowerPoint Presentation")

            if st.button("Generate PowerPoint Presentation"):
                pptx_bytes = create_presentation(df, ai_insights, sentiment, risks, recommendations)
                st.download_button("📥 Download PowerPoint", pptx_bytes, file_name="Financial_Report_Analysis.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    else:
        st.error("❌ No structured data found in the file. Please check the file and try again.")